"""
DocGen Render Engine
Supports: .docx (docxtpl), .pptx (python-pptx), PDF (weasyprint)

Required packages:
pip install docxtpl python-pptx weasyprint jinja2
"""

import os
import logging
from datetime import datetime
from pathlib import Path
from jinja2 import Environment

logger = logging.getLogger(__name__)

UPLOAD_FOLDER = os.environ.get('DOCGEN_UPLOAD_FOLDER', 'uploads/templates')
OUTPUT_FOLDER = os.environ.get('DOCGEN_OUTPUT_FOLDER', 'uploads/generated')

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


def render_docx(template_path: str, context: dict, output_path: str) -> str:
    """Fill a .docx template with Jinja2-style placeholders using docxtpl."""
    try:
        from docxtpl import DocxTemplate
        
        # Check if template file exists
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        logger.info(f"Loading DOCX template from: {template_path}")
        tpl = DocxTemplate(template_path)
        
        logger.info(f"Rendering with context keys: {list(context.keys())}")
        tpl.render(context)
        
        logger.info(f"Saving to: {output_path}")
        tpl.save(output_path)
        
        logger.info(f"DOCX rendered → {output_path}")
        return output_path
    except ImportError:
        raise ImportError("docxtpl not installed. Run: pip install docxtpl")
    except Exception as e:
        logger.error(f"DOCX rendering error: {e}")
        raise


def render_pptx(template_path: str, context: dict, output_path: str) -> str:
    """
    Fill a .pptx template.
    Replaces {{key}} placeholders in every text frame of every slide.
    """
    try:
        from pptx import Presentation
        prs = Presentation(template_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in context.items():
                            placeholder = '{{' + key + '}}'
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))

        prs.save(output_path)
        logger.info(f"PPTX rendered → {output_path}")
        return output_path
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")


def render_pdf(template_path: str, context: dict, output_path: str) -> str:
    """
    Render a Jinja2 HTML template to PDF using WeasyPrint.
    template_path: path to an .html file with Jinja2 syntax.
    """
    try:
        from weasyprint import HTML
        from jinja2 import Template

        with open(template_path, 'r', encoding='utf-8') as f:
            html_source = f.read()

        jinja_env = Environment()
        template = jinja_env.from_string(html_source)
        rendered_html = template.render(**context)

        HTML(string=rendered_html).write_pdf(output_path)
        logger.info(f"PDF rendered → {output_path}")
        return output_path
    except ImportError:
        raise ImportError("weasyprint not installed. Run: pip install weasyprint")


def docx_to_pdf(docx_path: str, output_path: str) -> str:
    """
    Convert a rendered .docx to PDF via LibreOffice (headless).
    Requires: apt-get install libreoffice
    Alternative if LibreOffice not available: use python-docx2pdf on Windows/Mac.
    
    NOTE: This function requires LibreOffice to be installed on the server.
    On Render Free Tier, LibreOffice is not available by default.
    """
    import subprocess
    
    # Check if libreoffice is available
    try:
        result = subprocess.run(['which', 'libreoffice'], capture_output=True, text=True, timeout=5)
        if result.returncode != 0:
            raise FileNotFoundError("LibreOffice not installed. PDF conversion not available. Please use DOCX output instead.")
    except Exception as e:
        raise FileNotFoundError(f"LibreOffice not available: {e}. Please use DOCX output instead.")

    out_dir = str(Path(output_path).parent)
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', out_dir, docx_path],
        capture_output=True, text=True, timeout=30
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    # LibreOffice names the output file after the input
    generated = Path(out_dir) / (Path(docx_path).stem + '.pdf')
    generated.rename(output_path)
    logger.info(f"DOCX→PDF → {output_path}")
    return output_path


def generate_document(template, record_data: dict, output_type: str = None) -> str:
    """
    Main entry point.
    template: DocTemplate model instance
    record_data: dict of CRM record fields
    output_type: 'pdf' | 'docx' | 'pptx' | None (uses template's native type)

    Returns: path to the generated file
    """
    context = _build_context(template, record_data)
    timestamp = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
    safe_name = template.name.replace(' ', '_').lower()
    output_type = output_type or template.file_type

    output_filename = f"{safe_name}_{record_data.get('id', 'unknown')}_{timestamp}.{output_type}"
    output_path = os.path.abspath(os.path.join(OUTPUT_FOLDER, output_filename))
    
    # Normalize template path
    template_path = os.path.abspath(template.file_path)
    
    logger.info(f"Generating document: template={template_path}, output={output_path}")

    if template.file_type == 'docx':
        if output_type == 'pdf':
            # Render docx first, then convert
            # Note: PDF conversion requires LibreOffice which may not be available on all platforms
            try:
                tmp_docx = output_path.replace('.pdf', '_tmp.docx')
                render_docx(template_path, context, tmp_docx)
                docx_to_pdf(tmp_docx, output_path)
                os.remove(tmp_docx)
            except FileNotFoundError as e:
                # LibreOffice not available - fall back to DOCX
                logger.warning(f"PDF conversion not available: {e}. Generating DOCX instead.")
                output_path = output_path.replace('.pdf', '.docx')
                render_docx(template_path, context, output_path)
        else:
            render_docx(template_path, context, output_path)

    elif template.file_type == 'pptx':
        render_pptx(template_path, context, output_path)

    elif template.file_type == 'html':
        render_pdf(template_path, context, output_path)

    else:
        raise ValueError(f"Unsupported template type: {template.file_type}")

    return output_path


def _build_context(template, record_data: dict) -> dict:
    """
    Map CRM record fields to template placeholders using template.field_map.
    field_map example: {"musteri_adi": "name", "toplam": "deal_value"}
    Falls back to using record_data directly if no field_map.
    """
    if not template.field_map:
        return record_data

    context = {}
    for placeholder, crm_field in template.field_map.items():
        context[placeholder] = record_data.get(crm_field, '')

    # Always include raw record data too, for direct access in templates
    context.update(record_data)
    return context
